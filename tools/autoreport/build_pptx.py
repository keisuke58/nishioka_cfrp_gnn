#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import csv
from pathlib import Path
from typing import Any, Dict, List, Optional


def _load_json(path: Path) -> Dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _try_import_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.text import PP_ALIGN  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
        return Presentation, PP_ALIGN, Inches, Pt
    except Exception as e:  # pragma: no cover
        raise SystemExit(
            "[FATAL] python-pptx が必要です。以下で入れてから実行してください:\n"
            "  /home/nishioka/miniconda3/envs/gnn_final_env/bin/python -m pip install -U python-pptx\n"
            f"  (import error: {e})"
        )


def _add_title(slide, title: str, subtitle: str, PP_ALIGN, Pt):
    title_box = slide.shapes.title
    if title_box:
        title_box.text = title
        p = title_box.text_frame.paragraphs[0]
        p.font.name = "Times New Roman"
        p.font.size = Pt(36)
    # subtitle placeholder
    if len(slide.placeholders) > 1:
        sub = slide.placeholders[1]
        sub.text = subtitle
        p = sub.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.name = "Times New Roman"
        p.font.size = Pt(18)


def _add_bullets(slide, header: str, bullets: List[str], Inches, Pt):
    left = Inches(0.8)
    top = Inches(1.3)
    width = Inches(12.0)
    height = Inches(5.0)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = header
    p0.font.name = "Times New Roman"
    p0.font.size = Pt(28)
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = "Times New Roman"
        p.font.size = Pt(20)


def _add_image_fit(slide, img_path: Path, Inches, x: float, y: float, max_w: float, max_h: float):
    if not img_path.exists():
        return
    # Fit image into a bounding box while keeping aspect ratio.
    try:
        from PIL import Image  # type: ignore

        with Image.open(img_path) as im:
            px_w, px_h = im.size
        if not px_w or not px_h:
            raise ValueError("invalid image size")

        img_aspect = float(px_w) / float(px_h)
        box_aspect = float(max_w) / float(max_h) if max_h else img_aspect
        if img_aspect >= box_aspect:
            w = float(max_w)
            h = w / img_aspect
        else:
            h = float(max_h)
            w = h * img_aspect

        x0 = float(x) + (float(max_w) - w) / 2.0
        y0 = float(y) + (float(max_h) - h) / 2.0
        slide.shapes.add_picture(str(img_path), Inches(x0), Inches(y0), width=Inches(w), height=Inches(h))
    except Exception:
        # Fallback: width-only (may overflow vertically, but avoids hard failure).
        slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), width=Inches(max_w))


def _find_asset(assets_dir: Path, patterns: List[str]) -> Optional[Path]:
    for pat in patterns:
        hits = sorted(assets_dir.glob(pat))
        if hits:
            return hits[0]
    return None

def _add_table(slide, title: str, columns: List[str], rows: List[List[str]], Inches, Pt):
    # Simple table helper
    slide.shapes.title.text = title
    if slide.shapes.title and slide.shapes.title.text_frame.paragraphs:
        p = slide.shapes.title.text_frame.paragraphs[0]
        p.font.name = "Times New Roman"
        p.font.size = Pt(32)
    left = Inches(0.6)
    top = Inches(1.4)
    width_in = 12.3
    height_in = 5.4
    width = Inches(width_in)
    height = Inches(height_in)
    n_rows = 1 + max(1, len(rows))
    n_cols = max(1, len(columns))
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table

    # Column widths (improve readability, avoid overflow)
    try:
        if n_cols == 2:
            table.columns[0].width = Inches(width_in * 0.78)
            table.columns[1].width = Inches(width_in * 0.22)
        else:
            each = width_in / float(n_cols)
            for j in range(n_cols):
                table.columns[j].width = Inches(each)
    except Exception:
        pass

    # Header
    for j, col in enumerate(columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.text_frame.word_wrap = True
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Times New Roman"
        p.font.size = Pt(14)

    # Body
    for i, r in enumerate(rows[: n_rows - 1], start=1):
        for j in range(n_cols):
            cell = table.cell(i, j)
            cell.text = str(r[j]) if j < len(r) else ""
            cell.text_frame.word_wrap = True
            for p in cell.text_frame.paragraphs:
                p.font.name = "Times New Roman"
                p.font.size = Pt(12)

    return table


def build_pptx(out_dir: Path) -> Path:
    Presentation, PP_ALIGN, Inches, Pt = _try_import_pptx()

    data = _load_json(out_dir / "report_data.json")
    assets_dir = out_dir / "assets"

    run_id = str(data.get("run_id", out_dir.name))
    profile = str(data.get("profile", ""))
    git_sha = str(data.get("git_short_sha", ""))

    best = ((data.get("results") or {}).get("best") or {})
    test = ((data.get("results") or {}).get("test") or {})
    ds = data.get("dataset") or {}
    art = data.get("artifacts") or {}
    stats = art.get("file_statistics_summary") or {}
    thr = art.get("threshold_decisions_summary") or {}
    train_sum = art.get("train_log_summary") or {}
    imbalance = train_sum.get("imbalance_summary") or {}
    cli_args = ((data.get("training") or {}).get("args") or [])

    def _extract(args_list: List[Any], key: str) -> Optional[str]:
        a = [str(x) for x in args_list]
        for i, tok in enumerate(a):
            if tok == f"--{key}":
                if i + 1 < len(a) and not a[i + 1].startswith("--"):
                    return a[i + 1]
                return "true"
        return None

    hp_lines = []
    for k in ["hidden_channels", "learning_rate", "weight_decay", "batch_size", "epochs", "patience", "dropout", "edge_drop_prob"]:
        v = _extract(cli_args, k)
        if v is not None:
            hp_lines.append(f"{k}: {v}")

    def _load_top_confusions(k: int = 10) -> List[List[str]]:
        p = assets_dir / "auto_top_confusions.csv"
        if not p.exists():
            return []
        rows: List[List[str]] = []
        try:
            with p.open("r", encoding="utf-8") as f:
                r = csv.DictReader(f)
                for i, row in enumerate(r):
                    if i >= k:
                        break
                    rows.append(
                        [
                            str(row.get("true_class", "")),
                            str(row.get("pred_class", "")),
                            str(row.get("count", "")),
                            str(row.get("row_rate", "")),
                        ]
                    )
        except Exception:
            return []
        return rows

    prs = Presentation()
    # Widescreen (16:9). Default is often 4:3, which causes our 12"+ layouts to overflow.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(
        slide,
        title=run_id,
        subtitle=f"profile={profile}  git={git_sha}",
        PP_ALIGN=PP_ALIGN,
        Pt=Pt,
    )

    def _set_title(slide, text: str):
        slide.shapes.title.text = text
        if slide.shapes.title and slide.shapes.title.text_frame.paragraphs:
            p = slide.shapes.title.text_frame.paragraphs[0]
            p.font.name = "Times New Roman"
            p.font.size = Pt(32)

    def _add_big_grid(slide, imgs: List[Path], title: str):
        _set_title(slide, title)
        # 2x2 grid (big)
        positions = [(0.7, 1.5), (6.9, 1.5), (0.7, 4.4), (6.9, 4.4)]
        for img, (x, y) in zip(imgs[:4], positions):
            _add_image_fit(slide, img, Inches, x=x, y=y, max_w=6.0, max_h=2.8)

    # Slide 2: Goal / Setting
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    _set_title(slide, "Goal & setup")
    _add_bullets(
        slide,
        header="Imbalanced 19-class defect classification",
        bullets=[
            "Primary metric: macro-F1 (accuracy can be dominated by the majority class)",
            f"Training script: {data.get('script','')}",
        ],
        Inches=Inches,
        Pt=Pt,
    )

    # Slide 3: Dataset
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Dataset")
    ds_rows = [
        ["type", str(ds.get("dataset_type", ""))],
        ["train/val/test", f"{ds.get('train_pairs','')}/{ds.get('val_pairs','')}/{ds.get('test_pairs','')}"],
        ["nodes/sample", str(ds.get("nodes_per_sample", ""))],
        ["#classes", str(ds.get("number_of_classes", ""))],
    ]
    _add_table(slide, "Dataset summary", ["item", "value"], ds_rows, Inches=Inches, Pt=Pt)

    # Slide 4: Class imbalance (quant)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Class imbalance (key facts)")
    if imbalance:
        _add_bullets(
            slide,
            header="Imbalance summary",
            bullets=[
                f"majority class: {imbalance.get('majority_class')} ({imbalance.get('majority_percent')}%)",
                f"minority class: {imbalance.get('minority_class')} ({imbalance.get('minority_samples')} samples, {imbalance.get('minority_percent')}%)",
                "Accuracy saturates with the majority class → optimize macro-F1 (minority recall)",
            ],
            Inches=Inches,
            Pt=Pt,
        )
    else:
        _add_bullets(slide, header="Imbalance summary", bullets=["Class distribution not found in train.log"], Inches=Inches, Pt=Pt)

    # Slide 5: Results (key metrics)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Results")
    res_rows = [
        ["best macro-F1 (val)", f"{best.get('macro_f1')} (epoch {best.get('epoch')})"],
        ["test macro-F1", str(test.get("macro_f1"))],
        ["test weighted-F1 / acc", f"{test.get('weighted_f1')} / {test.get('accuracy')}"],
        ["balanced-acc / MCC", f"{test.get('balanced_accuracy')} / {test.get('mcc')}"],
    ]
    _add_table(slide, "Key metrics", ["metric", "value"], res_rows, Inches=Inches, Pt=Pt)

    # Slide 6: Model (architecture overview)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Model overview (GAT, two-stage head)")
    _add_bullets(
        slide,
        header="Architecture",
        bullets=[
            "3× GATConv (heads=4) + BatchNorm + Dropout + Residual projections",
            "Two-stage head: (1) defect detection (2 classes) + (2) defect type classification (18 classes)",
            "Layer-constraint masking: physically impossible classes are masked with -inf logits",
            "Edge dropout is applied during training",
        ],
        Inches=Inches,
        Pt=Pt,
    )

    # Slide 6b: Two-stage probability composition
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Two-stage probability composition")
    _add_bullets(
        slide,
        header="Final 19-class probabilities",
        bullets=[
            "Detection: p_det(y=0|x), p_det(y>0|x)",
            "Classification (defect-only): p_cls(k|x, y>0), k=1..18",
            "Combine:",
            "  p(0|x) = p_det(y=0|x)",
            "  p(k|x) = p_det(y>0|x) * p_cls(k|x, y>0)",
        ],
        Inches=Inches,
        Pt=Pt,
    )

    # Slide 6c: Layer masking / physical constraints
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Layer masking (physical constraints)")
    _add_bullets(
        slide,
        header="Mask impossible classes per node layer",
        bullets=[
            "Allowed class set A(i) depends on layer of node i",
            "If k ∉ A(i): set logit to -inf (AMP-safe min value)",
            "Prevents physically impossible predictions and stabilizes training",
        ],
        Inches=Inches,
        Pt=Pt,
    )

    # Slide 6d: Training objective
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Training objective (two-stage loss)")
    _add_bullets(
        slide,
        header="Loss",
        bullets=[
            "Detection target: y_det = 1[y != 0] (2-class CE)",
            "Classification target: y_cls = y-1 for defect nodes (18-class CE)",
            "Total: L = λ_det * CE(det) + λ_cls * CE(cls on defect nodes)",
        ],
        Inches=Inches,
        Pt=Pt,
    )

    # Slide 6: Training hyperparameters (concise)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Training hyperparameters")
    if hp_lines:
        _add_bullets(slide, header="Hyperparameters", bullets=hp_lines[:10], Inches=Inches, Pt=Pt)
    else:
        _add_bullets(slide, header="Hyperparameters", bullets=["Failed to extract from CLI args"], Inches=Inches, Pt=Pt)

    # Slide 6e: Top confusions (from auto_top_confusions.csv)
    top_conf = _load_top_confusions(k=10)
    if top_conf:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        _set_title(slide, "Top confusions (from predictions)")
        _add_table(
            slide,
            "Most frequent confusions (top-10)",
            ["true", "pred", "count", "row_rate"],
            top_conf,
            Inches=Inches,
            Pt=Pt,
        )

    # Slide 7: Learning curves
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Learning curves")
    for name, x, y, w in [
        ("plot_combined_plot_20260115_224408.png", 0.6, 1.3, 12.0),
    ]:
        # prefer any combined_plot in assets
        combined = next(iter(sorted(assets_dir.glob("plot_combined_plot_*.png"))), None)
        if combined:
            _add_image_fit(slide, combined, Inches, x=0.6, y=1.3, max_w=12.0, max_h=5.9)
        else:
            # fallback: add loss/macro if present
            loss = next(iter(sorted(assets_dir.glob("plot_loss_plot_*.png"))), None)
            mf1 = next(iter(sorted(assets_dir.glob("plot_macro_f1_plot_*.png"))), None)
            if loss:
                _add_image_fit(slide, loss, Inches, x=0.6, y=1.3, max_w=6.0, max_h=5.9)
            if mf1:
                _add_image_fit(slide, mf1, Inches, x=6.8, y=1.3, max_w=6.0, max_h=5.9)

    # Slide 7b: Confusion matrix (auto from npy preferred)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Confusion matrix")
    cm = _find_asset(assets_dir, ["auto_confusion_counts.png", "truth_confusion_matrix*.png"])
    if cm:
        _add_image_fit(slide, cm, Inches, x=1.0, y=1.4, max_w=11.5, max_h=5.8)
    else:
        _add_bullets(slide, header="Confusion matrix", bullets=["Image not found"], Inches=Inches, Pt=Pt)

    # Slide 7c: Per-class F1
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Per-class F1")
    f1img = _find_asset(assets_dir, ["auto_per_class_f1.png", "truth_f1score_class*.png"])
    if f1img:
        _add_image_fit(slide, f1img, Inches, x=0.8, y=1.4, max_w=12.0, max_h=5.8)
    else:
        _add_bullets(slide, header="Per-class F1", bullets=["Image not found"], Inches=Inches, Pt=Pt)

    # Slide 7d: Confidence / Entropy (uncertainty)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Uncertainty (confidence / entropy)")
    conf = _find_asset(assets_dir, ["auto_confidence_hist.png"])
    ent = _find_asset(assets_dir, ["auto_entropy_hist.png"])
    if conf:
        _add_image_fit(slide, conf, Inches, x=0.8, y=1.6, max_w=5.8, max_h=5.4)
    if ent:
        _add_image_fit(slide, ent, Inches, x=6.6, y=1.6, max_w=5.8, max_h=5.4)
    if not conf and not ent:
        _add_bullets(slide, header="Uncertainty", bullets=["Confidence/entropy plots not found"], Inches=Inches, Pt=Pt)

    # Dataset properties (auto plots)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Dataset properties")
    a = _find_asset(assets_dir, ["dataset_defect_ndf_counts.png"])
    b = _find_asset(assets_dir, ["dataset_preddefectratio_hist.png"])
    if a:
        _add_image_fit(slide, a, Inches, x=0.8, y=1.6, max_w=5.8, max_h=5.4)
    if b:
        _add_image_fit(slide, b, Inches, x=6.6, y=1.6, max_w=5.8, max_h=5.4)
    if not a and not b:
        _add_bullets(slide, header="Dataset properties", bullets=["Dataset plots not found"], Inches=Inches, Pt=Pt)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Class distribution (counts / weights)")
    c1 = _find_asset(assets_dir, ["dataset_class_counts_log.png", "dataset_class_counts.png"])
    c2 = _find_asset(assets_dir, ["dataset_class_weights.png"])
    if c1:
        _add_image_fit(slide, c1, Inches, x=0.8, y=1.6, max_w=7.0, max_h=5.4)
    if c2:
        _add_image_fit(slide, c2, Inches, x=8.1, y=1.6, max_w=5.0, max_h=5.4)
    if not c1 and not c2:
        _add_bullets(slide, header="Class distribution", bullets=["Class distribution plots not found"], Inches=Inches, Pt=Pt)

    # MUST visualize good/bad examples (bigger images; more slides allowed)
    worst_imgs_all = sorted(assets_dir.glob("worst_*.png"))
    if not worst_imgs_all:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        _set_title(slide, "Worst examples (must visualize)")
        _add_bullets(slide, header="Worst examples", bullets=["No worst_*.png found"], Inches=Inches, Pt=Pt)
    else:
        for i in range(0, min(len(worst_imgs_all), 12), 4):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            _add_big_grid(slide, worst_imgs_all[i : i + 4], f"Worst examples (low per-file accuracy) [{i//4+1}]")

    best_imgs_all = sorted(assets_dir.glob("best_*.png"))
    if not best_imgs_all:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        _set_title(slide, "Best defect examples (must visualize)")
        _add_bullets(slide, header="Best examples", bullets=["No best_*.png found"], Inches=Inches, Pt=Pt)
    else:
        for i in range(0, min(len(best_imgs_all), 12), 4):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            _add_big_grid(slide, best_imgs_all[i : i + 4], f"Best defect examples (high per-file accuracy) [{i//4+1}]")

    # Slide 8: Failure analysis (tables)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    worst = stats.get("worst_accuracy", []) if isinstance(stats, dict) else []
    top_ratio = stats.get("top_pred_defect_ratio", []) if isinstance(stats, dict) else []
    rows = []
    for it in (worst[:6] if isinstance(worst, list) else []):
        rows.append([it.get("filename", ""), str(it.get("accuracy", ""))])
    if rows:
        _add_table(slide, "Failure analysis: lowest per-file accuracy (top)", ["filename", "accuracy"], rows, Inches=Inches, Pt=Pt)
    else:
        _add_bullets(slide, header="Failure analysis", bullets=["file_statistics_summary not found"], Inches=Inches, Pt=Pt)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    rows = []
    for it in (top_ratio[:6] if isinstance(top_ratio, list) else []):
        rows.append([it.get("filename", ""), str(it.get("pred_defect_ratio", ""))])
    if rows:
        _add_table(slide, "Failure analysis: highest PredDefectRatio (top)", ["filename", "pred_defect_ratio"], rows, Inches=Inches, Pt=Pt)
    else:
        _add_bullets(slide, header="Failure analysis", bullets=["top_pred_defect_ratio not found"], Inches=Inches, Pt=Pt)

    # Slide 10: Threshold decisions (quant)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tcounts = thr.get("threshold_counts") if isinstance(thr, dict) else None
    if isinstance(tcounts, dict) and tcounts:
        rows = [[k, str(v)] for k, v in sorted(tcounts.items())]
        _add_table(slide, "Thresholding: #files flagged as defect", ["column", "#files"], rows[:10], Inches=Inches, Pt=Pt)
    else:
        _add_bullets(slide, header="Thresholding", bullets=["file_decisions_thresholds not found"], Inches=Inches, Pt=Pt)

    # Slide 11: Spatial examples
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Spatial visualizations (selected)")
    imgs_all = sorted(assets_dir.glob("spatial_*.png"))
    # Use bigger images: 2x2 grid, and add more slides if needed
    if imgs_all:
        _add_big_grid(slide, imgs_all[:4], "Spatial visualizations (selected) [1]")
        for i in range(4, min(len(imgs_all), 12), 4):
            s2 = prs.slides.add_slide(prs.slide_layouts[5])
            _add_big_grid(s2, imgs_all[i : i + 4], f"Spatial visualizations (selected) [{i//4+1}]")
    else:
        _add_bullets(slide, header="Spatial visualizations", bullets=["No spatial_*.png found"], Inches=Inches, Pt=Pt)

    # Slide 12: Conclusion / next actions (quantified)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Conclusion & next actions")
    _add_bullets(
        slide,
        header="Takeaways",
        bullets=[
            f"Macro-F1 is the bottleneck: best(val)={best.get('macro_f1')}, test={test.get('macro_f1')}",
            "Next: improve minority-class recall (macro-F1 driven)",
            "Plan: (1) focal vs logit-adjust vs sampler (2) top confusions + qualitative review (3) threshold calibration",
        ],
        Inches=Inches,
        Pt=Pt,
    )

    out_path = out_dir / "slides.pptx"
    prs.save(str(out_path))
    return out_path


def main() -> int:
    ap = argparse.ArgumentParser(description="Build slides.pptx from report_data.json + assets/")
    ap.add_argument("--out", required=True, help="reports/<run>/ directory that contains report_data.json and assets/")
    args = ap.parse_args()

    out_dir = Path(args.out).expanduser().resolve()
    if not (out_dir / "report_data.json").exists():
        raise SystemExit(f"[FATAL] report_data.json not found in: {out_dir}")
    build_pptx(out_dir)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

